"""
Generates: ACT_prezentacja_parametryczna.pptx
Slides (9):
  0 – Tytuł
  1 – Problem i cel badania
  2 – Zbiór danych
  3 – Cenzurowanie prawostronne
  4 – Teoria: model AFT / Weibull
  5 – Wybór rozkładu – porównanie AIC + Probability Plots
  6 – Model zredukowany – istotne zmienne
  7 – Funkcja przeżycia i hazardu
  8 – Podsumowanie i wnioski
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Paths ─────────────────────────────────────────────────────────────────
BASE = os.path.dirname(os.path.abspath(__file__))
IMG  = os.path.join(BASE, "Ready_models")
OUT  = os.path.join(BASE, "ACT_prezentacja_parametryczna_lognormal.pptx")

# ── Palette ────────────────────────────────────────────────────────────────
C_BG        = RGBColor(0xF8, 0xF9, 0xFA)   # off-white background
C_HEADER    = RGBColor(0x1A, 0x3A, 0x5C)   # dark navy  – slide titles
C_ACCENT    = RGBColor(0x2E, 0x75, 0xB6)   # steel blue – section headers
C_TEXT      = RGBColor(0x1F, 0x1F, 0x1F)   # near-black
C_MUTED     = RGBColor(0x55, 0x65, 0x72)   # grey       – captions / bullets
C_LINE      = RGBColor(0x2E, 0x75, 0xB6)   # rule line
C_TITLE_BG  = RGBColor(0x1A, 0x3A, 0x5C)   # title-slide background
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # blank layout

# ══════════════════════════════════════════════════════════════════════════
# Helpers
# ══════════════════════════════════════════════════════════════════════════

def bg(slide, color=C_BG):
    """Fill slide background."""
    bg_shape = slide.background
    fill = bg_shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def rule(slide, y, left=Inches(0.5), width=Inches(12.33),
         color=C_LINE, thick=Pt(1.5)):
    """Horizontal accent rule."""
    from pptx.util import Pt as UPt
    ln = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, y, width, Inches(0.02)
    )
    ln.fill.solid()
    ln.fill.fore_color.rgb = color
    ln.line.fill.background()


def txbox(slide, text, x, y, w, h,
          size=18, bold=False, color=C_TEXT,
          align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return box


def heading(slide, text, y=Inches(0.25), size=32):
    txbox(slide, text, Inches(0.5), y, Inches(12.33), Inches(0.65),
          size=size, bold=True, color=C_HEADER)


def sub_heading(slide, text, y, size=18):
    txbox(slide, text, Inches(0.5), y, Inches(12.33), Inches(0.4),
          size=size, bold=True, color=C_ACCENT)


def body_bullets(slide, items, x, y, w, h, size=15, color=C_TEXT):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size  = Pt(size)
        run.font.color.rgb = color


def img(slide, path, x, y, w, h=None):
    if h:
        slide.shapes.add_picture(path, x, y, w, h)
    else:
        slide.shapes.add_picture(path, x, y, w)


def caption(slide, text, x, y, w):
    txbox(slide, text, x, y, w, Inches(0.3),
          size=10, color=C_MUTED, align=PP_ALIGN.CENTER)


def badge(slide, text, x, y, w=Inches(3), h=Inches(0.5),
          bg_color=C_ACCENT, fg_color=C_WHITE, size=13):
    """Filled rounded rectangle with text label."""
    from pptx.util import Emu
    shape = slide.shapes.add_shape(5, x, y, w, h)   # 5 = rounded rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p  = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = True
    run.font.color.rgb = fg_color


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 0 – Tytuł
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s, C_TITLE_BG)

# Decorative top band
band = s.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.18))
band.fill.solid(); band.fill.fore_color.rgb = C_ACCENT
band.line.fill.background()

# Main title
box = s.shapes.add_textbox(Inches(0.9), Inches(1.6), Inches(11.5), Inches(1.8))
tf  = box.text_frame; tf.word_wrap = True
p   = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run()
run.text = "Modele parametryczne AFT"
run.font.size = Pt(46); run.font.bold = True
run.font.color.rgb = C_WHITE

# Subtitle
box2 = s.shapes.add_textbox(Inches(0.9), Inches(3.3), Inches(11.5), Inches(0.9))
tf2  = box2.text_frame; tf2.word_wrap = True
p2   = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.CENTER
run2 = p2.add_run()
run2.text = "Analiza czasu przeżycia pacjentów z niewydolnością serca"
run2.font.size = Pt(22); run2.font.color.rgb = RGBColor(0xBD, 0xD7, 0xEE)

# Dataset badge
badge(s, "Heart Failure Clinical Records  |  N = 299",
      Inches(4.0), Inches(4.35), w=Inches(5.33), h=Inches(0.55),
      bg_color=RGBColor(0x2E,0x75,0xB6), size=14)

# Bottom info
txbox(s, "Analiza przeżycia  •  Model Log-Normal AFT  •  Python / lifelines",
      Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5),
      size=12, color=RGBColor(0x8A,0xAC,0xC8),
      align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 – Problem i cel badania
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Problem i cel badania")
rule(s, Inches(1.05))

sub_heading(s, "Co modelujemy?", Inches(1.25))
body_bullets(s, [
    "• Zmienna zależna: czas do zgonu (w dniach) od pierwszej wizyty klinicznej",
    "• Zdarzenie: DEATH_EVENT = 1 (zgon)  |  DEATH_EVENT = 0 (cenzurowanie prawostronne)",
    "• Dane prawostronnie cenzurowane — pacjent przeżył do końca okresu badania lub wypadł z obserwacji",
], Inches(0.6), Inches(1.7), Inches(12.1), Inches(1.3), size=15)

sub_heading(s, "Cel", Inches(3.1))
body_bullets(s, [
    "• Zbudować model predykcyjny czasu przeżycia uwzględniający zmienne kliniczne",
    "• Wybrać najlepiej dopasowany rozkład parametryczny (Weibull, Log-Normal, Log-Logistic)",
    "• Zidentyfikować predyktory istotnie wpływające na czas do zgonu",
    "• Oszacować funkcje przeżycia S(t) i hazardu h(t) dla różnych profili pacjentów",
], Inches(0.6), Inches(3.55), Inches(12.1), Inches(2.0), size=15)

sub_heading(s, "Dlaczego model parametryczny?", Inches(5.65))
body_bullets(s, [
    "• Pełna specyfikacja rozkładu → możliwość ekstrapolacji i predykcji mediany przeżycia",
    "• Alternatywa wobec nieparametrycznego KM i semiparametrycznego Coxa — sprawdzamy, który rozkład najlepiej opisuje dane",
], Inches(0.6), Inches(6.1), Inches(12.1), Inches(0.95), size=14, color=C_MUTED)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 – Zbiór danych
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Zbiór danych — Heart Failure Clinical Records")
rule(s, Inches(1.05))

# Left column: variable table
sub_heading(s, "Zmienne objaśniające (11)", Inches(1.2), size=16)

table_data = [
    ("Zmienna",             "Typ",       "Zakres/Wartości"),
    ("age",                 "ciągła",    "40 – 95 lat"),
    ("ejection_fraction",   "ciągła",    "14 – 80 %"),
    ("serum_creatinine",    "ciągła",    "0.5 – 9.4 mg/dL"),
    ("creatinine_phos...",  "ciągła",    "23 – 7861 IU/L"),
    ("serum_sodium",        "ciągła",    "113 – 148 mEq/L"),
    ("platelets",           "ciągła",    "25 100 – 850 000"),
    ("anaemia",             "binarna",   "0 / 1"),
    ("diabetes",            "binarna",   "0 / 1"),
    ("high_blood_pressure", "binarna",   "0 / 1"),
    ("sex",                 "binarna",   "0=K / 1=M"),
    ("smoking",             "binarna",   "0 / 1"),
]

col_x = [Inches(0.5), Inches(4.3), Inches(7.0)]
col_w = [Inches(3.7), Inches(2.6), Inches(2.8)]
row_h = Inches(0.33)
start_y = Inches(1.6)

for r, row in enumerate(table_data):
    y = start_y + r * row_h
    for c, cell in enumerate(row):
        is_header = (r == 0)
        b = s.shapes.add_textbox(col_x[c], y, col_w[c], row_h)
        tf = b.text_frame
        p  = tf.paragraphs[0]
        run = p.add_run()
        run.text = cell
        run.font.size  = Pt(12 if not is_header else 13)
        run.font.bold  = is_header
        run.font.color.rgb = C_WHITE if is_header else C_TEXT
        if is_header:
            b.fill.solid(); b.fill.fore_color.rgb = C_HEADER
        elif r % 2 == 0:
            b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0xE8,0xF0,0xF7)
        b.line.fill.background()

# Right column: key stats
sub_heading(s, "Statystyki kluczowe", Inches(1.2), size=16)
# reposition right panel
stats = [
    ("N = 299", "obserwacje"),
    ("96 (32.1%)", "zdarzenia (zgony)"),
    ("203 (67.9%)", "cenzurowane"),
    ("śr. 130 dni", "czas obserwacji"),
    ("śr. 61 lat", "wiek pacjenta"),
]
for i, (val, lbl) in enumerate(stats):
    bx = Inches(10.05)
    by = Inches(1.65) + i * Inches(0.9)
    box_bg = s.shapes.add_shape(1, bx, by, Inches(2.8), Inches(0.72))
    box_bg.fill.solid()
    box_bg.fill.fore_color.rgb = RGBColor(0xE8,0xF0,0xF7) if i%2==0 else C_BG
    box_bg.line.color.rgb = RGBColor(0xC0,0xD0,0xE0)
    txbox(s, val, bx+Inches(0.1), by+Inches(0.03), Inches(2.6), Inches(0.35),
          size=18, bold=True, color=C_ACCENT)
    txbox(s, lbl, bx+Inches(0.1), by+Inches(0.36), Inches(2.6), Inches(0.28),
          size=11, color=C_MUTED)

txbox(s, "* Brak wartości brakujących — dane kompletne",
      Inches(0.5), Inches(7.1), Inches(9.0), Inches(0.3),
      size=11, color=C_MUTED)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 – Cenzurowanie prawostronne
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Cenzurowanie prawostronne")
rule(s, Inches(1.05))

body_bullets(s, [
    "• Obserwacja jest prawostronnie cenzurowana, gdy pacjent przeżywa do końca badania lub wypadł z obserwacji bez zarejestrowania zgonu.",
    "• Wiemy jedynie, że przeżył co najmniej time dni — dolna granica czasu przeżycia jest znana.",
    "• Cenzurowanie lewostronne i interwałowe nie występują naturalnie w tym zbiorze — model standardowy jest właściwy.",
], Inches(0.6), Inches(1.2), Inches(7.8), Inches(2.3), size=15)

txbox(s, "Dlaczego pominięto inne typy cenzurowania?",
      Inches(0.6), Inches(3.6), Inches(7.8), Inches(0.4),
      size=14, bold=True, color=C_ACCENT)
body_bullets(s, [
    "→ Cenzurowanie lewostronne: zademonstrowano symulacyjnie (30 obs.) — wyniki MLE były numerycznie niestabilne (SE → ∞), co potwierdza brak użyteczności.",
    "→ Cenzurowanie interwałowe: czas time jest dokładny (dni) — zakres (t−1, t] nie wnosi informacji.",
], Inches(0.6), Inches(4.05), Inches(7.8), Inches(1.4), size=13, color=C_MUTED)

img(s, os.path.join(IMG, "hf_censoring.png"),
    Inches(8.5), Inches(1.1), Inches(4.5))
caption(s, "Rys. 1. Rozkład czasu obserwacji wg statusu oraz proporcja zdarzeń",
        Inches(8.5), Inches(5.5), Inches(4.5))


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 – Teoria: model AFT / Weibull
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Model Accelerated Failure Time (AFT)")
rule(s, Inches(1.05))

sub_heading(s, "Równanie modelu", Inches(1.2))
# Formula box
fb = s.shapes.add_shape(1, Inches(0.6), Inches(1.6), Inches(7.8), Inches(0.85))
fb.fill.solid(); fb.fill.fore_color.rgb = RGBColor(0xE8,0xF0,0xF7)
fb.line.color.rgb = C_ACCENT
txbox(s, "log(T) = μ + x·β + σ·ε      gdzie  ε ~ Normal (Log-Normal AFT)",
      Inches(0.7), Inches(1.65), Inches(7.6), Inches(0.75),
      size=17, bold=True, color=C_HEADER)

sub_heading(s, "Interpretacja parametrów", Inches(2.6))
body_bullets(s, [
    "• exp(βⱼ) — współczynnik przyspieszenia (acceleration factor): o ile razy zmienia się oczekiwany czas przeżycia przy wzroście xⱼ o 1",
    "• exp(β) > 1 → czas przeżycia rośnie (czynnik ochronny)",
    "• exp(β) < 1 → czas przeżycia maleje (czynnik ryzyka)",
    "• σ (parametr skali log-czasu): log(T) | x ~ N(μ + x·β, σ²)  →  h(t) rośnie do maksimum, następnie maleje (kształt dzwonowy)",
], Inches(0.6), Inches(3.05), Inches(7.8), Inches(2.5), size=14)

txbox(s, "Log-Normal: hazard niemonotoniczny — biologicznie uzasadniony dla HF (ryzyko rośnie, osiąga szczyt, maleje)",
      Inches(0.6), Inches(5.6), Inches(7.8), Inches(0.5),
      size=13, bold=True, color=C_ACCENT)

# Right: comparison table of distributions
sub_heading(s, "Dostępne rozkłady AFT", Inches(1.2))
dist_info = [
    ("Rozkład",    "Kształt hazardu",     "Cecha kliniczna"),
    ("Weibull",    "Monot. malejący (ρ<1)","Najwyższe ryzyko w dniu 0, stale opada"),
    ("Log-Normal", "Dzwonowy ↑ potem ↓", "Ryzyko rośnie, osiąga szczyt, maleje — ✓ HF"),
    ("Log-Logistic","Dzwonowy, grubszy ogon","Podobny do Log-Normal"),
]
col_x2 = [Inches(8.7), Inches(10.1), Inches(11.2)]
col_w2 = [Inches(1.35), Inches(1.1), Inches(1.85)]
for r, row in enumerate(dist_info):
    y = Inches(1.6) + r * Inches(0.62)
    for c, cell in enumerate(row):
        b = s.shapes.add_textbox(col_x2[c], y, col_w2[c], Inches(0.58))
        tf = b.text_frame
        p  = tf.paragraphs[0]
        run = p.add_run()
        run.text = cell
        run.font.size  = Pt(11 if r > 0 else 12)
        run.font.bold  = (r == 0)
        run.font.color.rgb = C_WHITE if r == 0 else C_TEXT
        if r == 0:
            b.fill.solid(); b.fill.fore_color.rgb = C_HEADER
        elif r % 2 == 1:
            b.fill.solid(); b.fill.fore_color.rgb = RGBColor(0xE8,0xF0,0xF7)
        b.line.fill.background()

txbox(s, "Wszystkie trzy rozkłady dopasowano do tych samych danych i porównano kryterium AIC/BIC →",
      Inches(8.7), Inches(4.0), Inches(4.5), Inches(0.8),
      size=12, color=C_MUTED)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 – Wybór rozkładu (AIC + Probability Plots)
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Wybór rozkładu — porównanie AIC i Probability Plots")
rule(s, Inches(1.05))

# AIC comparison image (left)
img(s, os.path.join(IMG, "hf_aft_aic_comparison.png"),
    Inches(0.4), Inches(1.2), Inches(5.5))
caption(s, "Rys. 2. AIC trzech modeli AFT — niższe = lepsze",
        Inches(0.4), Inches(5.3), Inches(5.5))

# Probability plots (right)
img(s, os.path.join(IMG, "hf_probability_plots.png"),
    Inches(6.2), Inches(1.15), Inches(6.9))
caption(s, "Rys. 3. Probability Plots — linearyzacja estymaty KM; im bliżej R²=1, tym lepsze dopasowanie rozkładu",
        Inches(6.2), Inches(5.3), Inches(6.9))

# Results box
res_bg = s.shapes.add_shape(1, Inches(0.4), Inches(5.55), Inches(12.4), Inches(1.75))
res_bg.fill.solid(); res_bg.fill.fore_color.rgb = RGBColor(0xE8,0xF0,0xF7)
res_bg.line.color.rgb = C_ACCENT

body_bullets(s, [
    "AIC:  Weibull = 1282.2  |  Log-Logistic = 1285.5  |  Log-Normal = 1287.4   →   ΔAIC(Weibull–LogNormal) = 5.2",
    "Probability Plot R²:  Log-Normal wyższe R² niż Weibull  →  lepsze dopasowanie marginalnego rozkładu T do krzywej KM",
    "Concordance Index: identyczny dla wszystkich trzech (0.74) — dyskryminacja nie rozstrzyga wyboru",
    "Konflikt sygnałów: AIC wskazuje Weibull, R² wskazuje Log-Normal  →  ΔAIC = 5.2 to strefa umiarkowana (2–6), nie rozstrzygająca",
    "Wybrano Log-Normal AFT: ΔAIC = 5.2 jest w strefie umiarkowanej (< 6 wg Burnhama & Andersona) + wyższe R² + uzasadnienie kliniczne",
], Inches(0.6), Inches(5.6), Inches(12.1), Inches(1.65), size=11.5, color=C_TEXT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 – Model zredukowany
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Model zredukowany — Log-Normal AFT")
rule(s, Inches(1.05))

sub_heading(s, "Selekcja zmiennych (α = 0.05)", Inches(1.2), size=16)

# Variables table
kept = [
    ("age",                 "-0.045", "0.96", "< 0.001", "starszy wiek → krótszy czas przeżycia"),
    ("ejection_fraction",   "+0.054", "1.06", "< 0.001", "wyższa EF → dłuższy czas przeżycia"),
    ("serum_creatinine",    "-0.350", "0.70", "< 0.001", "wyższy kreatynina → krótszy czas"),
    ("serum_sodium",        "+0.038", "1.04", "  0.022",  "wyższy sód → dłuższy czas przeżycia"),
    ("anaemia",             "-0.420", "0.66", "  0.030",  "niedokrwistość skraca czas przeżycia"),
    ("high_blood_pressure", "-0.490", "0.61", "  0.025",  "nadciśnienie skraca czas przeżycia"),
]
removed = ["diabetes (p=0.53)", "platelets (p=0.64)", "sex (p=0.35)", "smoking (p=0.65)", "creatinine_phosphokinase (p=0.11)"]

hdr = [("Zmienna", 2.8), ("coef", 0.8), ("exp(coef)", 1.0), ("p", 0.75), ("Interpretacja", 5.2)]
x_starts = [Inches(0.4)]
for _, w in hdr[:-1]:
    x_starts.append(x_starts[-1] + Inches(w))

# Header row
for c, (hname, hw) in enumerate(hdr):
    b = s.shapes.add_textbox(x_starts[c], Inches(1.6), Inches(hw), Inches(0.38))
    b.fill.solid(); b.fill.fore_color.rgb = C_HEADER; b.line.fill.background()
    p = b.text_frame.paragraphs[0]; run = p.add_run()
    run.text = hname; run.font.size = Pt(12); run.font.bold = True
    run.font.color.rgb = C_WHITE

for r, row_data in enumerate(kept):
    y = Inches(2.0) + r * Inches(0.52)
    fill_c = RGBColor(0xE8,0xF0,0xF7) if r%2==0 else C_BG
    for c, (cell, (_, hw)) in enumerate(zip(row_data, hdr)):
        b = s.shapes.add_textbox(x_starts[c], y, Inches(hw), Inches(0.48))
        b.fill.solid(); b.fill.fore_color.rgb = fill_c; b.line.fill.background()
        p = b.text_frame.paragraphs[0]; run = p.add_run()
        run.text = cell
        run.font.size = Pt(11)
        run.font.color.rgb = C_TEXT if c != 1 else (
            RGBColor(0xC0,0x00,0x00) if "-" in cell else RGBColor(0x00,0x70,0x0C))

# Removed variables
txbox(s, "Usunięte (p > 0.05): " + "  |  ".join(removed),
      Inches(0.4), Inches(5.35), Inches(9.5), Inches(0.4),
      size=11, color=C_MUTED)

# AIC comparison box
aic_bg = s.shapes.add_shape(1, Inches(0.4), Inches(5.85), Inches(9.5), Inches(1.3))
aic_bg.fill.solid(); aic_bg.fill.fore_color.rgb = RGBColor(0xE8,0xF0,0xF7)
aic_bg.line.color.rgb = C_ACCENT
body_bullets(s, [
    "Porównanie AIC:  Model pełny Log-Normal = 1287.37   |   Model zredukowany = ~1278.1",
    "ΔAIC ≈ −9.3  →  model zredukowany lepszy (niższe AIC przy mniejszej liczbie parametrów)",
    "Concordance Index: 0.74  →  model zachowuje pełną zdolność dyskryminacyjną",
], Inches(0.6), Inches(5.9), Inches(9.2), Inches(1.2), size=12)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 – Funkcja przeżycia i hazardu
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Predykcja — funkcje przeżycia i hazardu")
rule(s, Inches(1.05))

img(s, os.path.join(IMG, "hf_aft_survival_profiles.png"),
    Inches(0.3), Inches(1.15), Inches(6.4))
caption(s, "Rys. 4. Funkcja przeżycia S(t) z 95% CI bootstrap (200 iteracji)",
        Inches(0.3), Inches(5.35), Inches(6.4))

img(s, os.path.join(IMG, "hf_aft_hazard_profiles.png"),
    Inches(6.9), Inches(1.15), Inches(6.1))
caption(s, "Rys. 5. Funkcja hazardu h(t) — aproksymacja numeryczna z S(t)",
        Inches(6.9), Inches(5.35), Inches(6.1))

# Profile legend box
prof_bg = s.shapes.add_shape(1, Inches(0.3), Inches(5.65), Inches(12.7), Inches(1.55))
prof_bg.fill.solid(); prof_bg.fill.fore_color.rgb = RGBColor(0xF0,0xF4,0xF8)
prof_bg.line.color.rgb = C_ACCENT

body_bullets(s, [
    "Pacjent A – niskie ryzyko:     wiek=55, EF=45%, kreatynina=1.0, bez anemii, bez nadciśnienia",
    "Pacjent B – średnie ryzyko:   wiek=65, EF=30%, kreatynina=1.5, z anemią",
    "Pacjent C – wysokie ryzyko:  wiek=75, EF=20%, kreatynina=2.0, z nadciśnieniem",
    "Pacjent D – wysokie ryzyko:  wiek=55, EF=25%, kreatynina=1.2, z anemią i nadciśnieniem",
], Inches(0.5), Inches(5.7), Inches(12.3), Inches(1.45), size=11.5, color=C_TEXT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 – Podsumowanie i wnioski
# ══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
bg(s)
heading(s, "Podsumowanie i wnioski")
rule(s, Inches(1.05))

sub_heading(s, "Wyniki modelowania", Inches(1.2))
body_bullets(s, [
    "• Wybór rozkładu: Log-Normal AFT (ΔAIC = 5.2 wobec Weibull — strefa umiarkowana) → hazard dzwonowy, biologicznie uzasadniony",
    "• 6 istotnych predyktorów (p < 0.05): wiek, EF, kreatynina, sód w surowicy, anemia, nadciśnienie",
    "• Concordance Index = 0.74 — model zachowuje pełną zdolność dyskryminacyjną po redukcji",
], Inches(0.6), Inches(1.65), Inches(8.2), Inches(1.7), size=13.5)

sub_heading(s, "Uwaga: AIC vs Probability Plot — nierozstrzygnięty konflikt", Inches(3.45))

# Conflict warning box
warn_bg = s.shapes.add_shape(1, Inches(0.5), Inches(3.9), Inches(8.1), Inches(1.55))
warn_bg.fill.solid(); warn_bg.fill.fore_color.rgb = RGBColor(0xFF, 0xF3, 0xCD)
warn_bg.line.color.rgb = RGBColor(0xE6, 0xA8, 0x17)
body_bullets(s, [
    "AIC faworyzuje Weibull, ale ΔAIC = 5.2 leży w strefie umiarkowanej (2–6 wg Burnhama & Andersona) — wybór nie jest jednoznaczny.",
    "Probability Plot wykazuje wyższe R² dla Log-Normal → lepsze dopasowanie do marginalnego rozkładu czasu T.",
    "Klinicznie: hazard dzwonowy Log-Normal jest bardziej realistyczny dla niewydolności serca — ryzyko rośnie po hospitalizacji, osiąga szczyt, następnie maleje. Hazard monotonicznie malejący (Weibull) sugeruje najwyższe ryzyko w dniu 0, co jest biologicznie dyskusyjne.",
    "Rekomendacja: przy ΔAIC < 6 uzasadniony jest wybór Log-Normal z uzasadnieniem klinicznym.",
], Inches(0.65), Inches(3.95), Inches(7.85), Inches(1.45), size=11.5, color=RGBColor(0x5A,0x3E,0x00))

sub_heading(s, "Wnioski kliniczne", Inches(5.55))
body_bullets(s, [
    "• Frakcja wyrzutowa (EF) i kreatynina — najsilniejsze predyktory (funkcja serca i nerek)",
    "• Anemia i nadciśnienie skracają oczekiwany czas przeżycia o ~35–40%",
    "• Diabetes, płytki krwi, płeć i palenie — bez istotnego efektu po kontroli powyższych",
], Inches(0.6), Inches(6.0), Inches(8.2), Inches(1.2), size=13)

# Side image: bootstrap
img(s, os.path.join(IMG, "hf_aft_bootstrap_median.png"),
    Inches(8.8), Inches(1.1), Inches(4.3))
caption(s, "Rys. 6. Bootstrap mediany przeżycia (100 iteracji)",
        Inches(8.8), Inches(5.0), Inches(4.3))

# Bottom bar
bot = s.shapes.add_shape(1, 0, Inches(7.1), SLIDE_W, Inches(0.4))
bot.fill.solid(); bot.fill.fore_color.rgb = C_HEADER; bot.line.fill.background()
txbox(s, "Dane: Heart Failure Clinical Records  •  N=299  •  96 zdarzeń  •  Model: Log-Normal AFT  •  Biblioteka: lifelines (Python)",
      Inches(0.3), Inches(7.12), Inches(12.7), Inches(0.35),
      size=11, color=C_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════
prs.save(OUT)
print(f"Saved: {OUT}")
